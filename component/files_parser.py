# -*- coding: utf-8 -*-
"""
文件处理主模块（整合版 - 包含图标系统）
包含：文本提取、智能分块、MinIO存储、文件验证、文件图标系统
"""

import os
import io
import logging
import re
import json
import traceback
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional, Union

from minio import Minio
from minio.error import S3Error
from .memory_sqlite import save_file_mapping
from .logger import LoggerManager
from .config import Config

# 初始化日志
logger = LoggerManager.get_logger(__name__)

# 初始化 MinIO 客户端
minio_client = Minio(
    Config.MINIO_ENDPOINT,
    access_key=Config.MINIO_ACCESS_KEY,
    secret_key=Config.MINIO_SECRET_KEY,
    secure=Config.MINIO_SECURE
)

# 导入 rag_course 模块的函数（避免循环导入）
from .rag_course import add_documents_to_store_from_minio

# ==================== 文件图标映射配置 ====================

FILE_TYPE_ICONS = {
    # 文档类图标（仅保留指定的）
    'pdf': 'PDF',
    'docx': 'DOCX',
    'pptx': 'PPT',
    'txt': 'TXT',
    'html': 'HTML',
    'ipynb': 'IPYNB',
    
    # 默认图标
    'default': 'DEFAULT'
}

FILE_TYPE_COLORS = {
    'pdf': '#dc2626',      # 红色
    'docx': '#2563eb',     # 蓝色
    'pptx': '#d97706',     # 橙色
    'txt': '#059669',      # 绿色
    'html': '#7c3aed',     # 紫色
    'ipynb': '#7c3aed',     # 紫色
    'default': '#6b7280'    # 灰色
}

# 合并基础图标映射（不再使用扩展映射）
ALL_FILE_TYPE_ICONS = FILE_TYPE_ICONS.copy()

# ==================== 文件图标获取函数 ====================

def get_file_type_icon(file_extension: str, default_icon: str = '📄') -> str:
    """
    根据文件扩展名获取对应的图标（精简版）
    
    Args:
        file_extension: 文件扩展名（如 'pdf', 'docx'）
        default_icon: 默认图标
        
    Returns:
        对应的文件图标
    """
    file_extension = file_extension.lower().lstrip('.').strip()
    
    # 只支持指定的文件类型
    supported_extensions = ['pdf', 'docx', 'pptx', 'txt', 'html', 'ipynb']
    
    if file_extension in supported_extensions:
        return ALL_FILE_TYPE_ICONS.get(file_extension, default_icon)
    else:
        # 不支持的文件类型返回默认图标
        return default_icon

def get_file_type_color(file_extension: str, default_color: str = '#6b7280') -> str:
    """
    根据文件扩展名获取对应的颜色（精简版）
    
    Args:
        file_extension: 文件扩展名（如 'pdf', 'docx'）
        default_color: 默认颜色
        
    Returns:
        对应的文件颜色
    """
    file_extension = file_extension.lower().lstrip('.').strip()
    
    # 只支持指定的文件类型
    supported_extensions = ['pdf', 'docx', 'pptx', 'txt', 'html', 'ipynb']
    
    if file_extension in supported_extensions:
        return FILE_TYPE_COLORS.get(file_extension, default_color)
    else:
        # 不支持的文件类型返回默认颜色
        return default_color

# ==================== 文件卡片HTML生成函数 ====================

def generate_file_card_html(file_name: str, file_size: int, file_extension: str, minio_path: str = "") -> str:
    """
    生成文件卡片HTML（优化版 - 不显示 minio_path）
    
    Args:
        file_name: 文件名
        file_size: 文件大小（字节）
        file_extension: 文件扩展名
        minio_path: MinIO 存储路径（可选，不显示）
        
    Returns:
        文件卡片HTML字符串
    """
    # 获取文件图标和颜色
    file_icon = get_file_type_icon(file_extension)
    file_color = get_file_type_color(file_extension)
    
    # 格式化文件大小
    if file_size < 1024:
        size_str = f"{file_size} B"
    elif file_size < 1024 * 1024:
        size_str = f"{file_size / 1024:.2f} KB"
    else:
        size_str = f"{file_size / (1024 * 1024):.2f} MB"
    
    # 生成文件卡片HTML（不显示 minio_path）
    file_card_html = f"""
    <div class="file-card" style="border-left: 4px solid {file_color};">
        <div class="file-card-icon" style="color: {file_color};">
            {file_icon}
        </div>
        <div class="file-card-info">
            <div class="file-card-name" style="color: {file_color};">
                {file_name}
            </div>
            <div class="file-card-meta">
                <span class="file-size">{size_str}</span>
                <span class="file-type">{file_extension.upper()}</span>
            </div>
        </div>
        <div class="file-card-status" style="margin-left: auto;">
            <span class="status-icon">✅</span>
        </div>
    </div>
    """
    
    return file_card_html

def generate_file_status_card(file_name: str, status: str, message: str, file_extension: str = None) -> str:
    """
    生成文件状态卡片HTML（优化版 - 不显示 minio_path）
    
    Args:
        file_name: 文件名
        status: 状态类型 ('success', 'failed', 'processing', 'warning')
        message: 状态消息
        file_extension: 文件扩展名（可选）
        
    Returns:
        状态卡片HTML字符串
    """
    # 获取状态图标和颜色
    status_config = {
        'success': {'icon': '✅', 'color': '#10b981', 'bg_color': '#ecfdf5'},
        'failed': {'icon': '❌', 'color': '#ef4444', 'bg_color': '#fef2f2'},
        'processing': {'icon': '⏳', 'color': '#3b82f6', 'bg_color': '#eff6ff'},
        'warning': {'icon': '⚠️', 'color': '#f59e0b', 'bg_color': '#fffbeb'},
        'info': {'icon': 'ℹ️', 'color': '#6b7280', 'bg_color': '#f3f4f6'}
    }
    
    config = status_config.get(status, status_config['info'])
    status_icon = config['icon']
    status_color = config['color']
    status_bg_color = config['bg_color']
    
    # 获取文件图标
    file_icon = get_file_type_icon(file_extension) if file_extension else '📄'
    
    # 生成状态卡片HTML（不显示 minio_path）
    status_card_html = f"""
    <div class="file-status-card" style="background-color: {status_bg_color}; border-left: 4px solid {status_color};">
        <div class="status-icon" style="color: {status_color};">
            {status_icon}
        </div>
        <div class="status-info">
            <div class="file-name" style="color: {status_color};">
                {file_name}
            </div>
            <div class="status-message" style="color: {status_color};">
                {message}
            </div>
        </div>
        <div class="file-icon" style="margin-left: auto; font-size: 24px;">
            {file_icon}
        </div>
    </div>
    """
    
    return status_card_html

# ==================== 文本提取模块 ====================

def extract_text_from_file(file_path: str, file_extension: Optional[str] = None) -> str:
    """
    从文件中提取文本内容
    
    Args:
        file_path: 文件路径
        file_extension: 文件扩展名（可选）
        
    Returns:
        提取的文本内容
    """
    try:
        # 检查文件是否存在
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")
        
        # 检查文件是否可读
        if not os.access(file_path, os.R_OK):
            raise PermissionError(f"没有权限读取文件: {file_path}")
        
        # 获取文件大小
        file_size = os.path.getsize(file_path)
        logger.info(f"文件大小: {file_size} 字节")
        
        if file_size == 0:
            raise ValueError(f"文件为空: {file_path}")
        
        # 如果没有提供文件扩展名，从文件路径中提取
        if not file_extension:
            file_extension = os.path.splitext(file_path)[1].lower().lstrip('.')
        
        # 根据文件扩展名选择合适的提取方法
        if file_extension == 'pdf':
            return extract_text_from_pdf(file_path)
        elif file_extension == 'docx':
            return extract_text_from_docx(file_path)
        elif file_extension == 'pptx':
            return extract_text_from_pptx(file_path)
        elif file_extension == 'txt':
            return extract_text_from_txt(file_path)
        elif file_extension == 'html':
            return extract_text_from_html(file_path)
        elif file_extension == 'ipynb':
            return extract_text_from_ipynb(file_path)
        else:
            logger.warning(f"不支持的文件类型: {file_extension}")
            return ""
    except Exception as e:
        logger.error(f"提取文件内容失败: {e}", exc_info=True)
        return ""

def extract_text_from_pdf(file_path: str) -> str:
    """从PDF文件中提取文本"""
    text = ""
    try:
        logger.info(f"开始提取PDF文件内容: {file_path}")
        
        # 优先使用 pypdf
        try:
            import pypdf
            logger.info("使用 pypdf 库提取PDF文本")
            with open(file_path, 'rb') as file:
                reader = pypdf.PdfReader(file)
                page_count = len(reader.pages)
                logger.info(f"PDF文件共有 {page_count} 页")
                
                # 检查PDF是否加密
                if reader.is_encrypted:
                    logger.warning("PDF文件已加密，尝试解密...")
                    try:
                        reader.decrypt('')
                        logger.info("PDF解密成功")
                    except Exception as e:
                        logger.error(f"PDF解密失败: {str(e)}")
                        raise ValueError("PDF文件已加密，无法提取文本")
                
                for page_num, page in enumerate(reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                        else:
                            logger.warning(f"第 {page_num+1} 页未能提取到文本")
                    except Exception as e:
                        logger.error(f"提取第 {page_num+1} 页文本失败: {str(e)}")
                        continue
        except ImportError:
            logger.warning("pypdf 库未安装，尝试使用 PyPDF2")
            try:
                import PyPDF2
                logger.info("使用 PyPDF2 库提取PDF文本")
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    page_count = pdf_reader.getNumPages()
                    logger.info(f"PDF文件共有 {page_count} 页")
                    
                    # 检查PDF是否加密
                    if pdf_reader.isEncrypted:
                        logger.warning("PDF文件已加密，尝试解密...")
                        try:
                            pdf_reader.decrypt('')
                            logger.info("PDF解密成功")
                        except Exception as e:
                            logger.error(f"PDF解密失败: {str(e)}")
                            raise ValueError("PDF文件已加密，无法提取文本")
                    
                    for page_num in range(page_count):
                        try:
                            page = pdf_reader.getPage(page_num)
                            page_text = page.extractText()
                            if page_text:
                                text += page_text + "\n"
                            else:
                                logger.warning(f"第 {page_num+1} 页未能提取到文本")
                        except Exception as e:
                            logger.error(f"提取第 {page_num+1} 页文本失败: {str(e)}")
                            continue
            except ImportError:
                logger.warning("PyPDF2 库也未安装，尝试使用 PyMuPDF")
                try:
                    import fitz  # PyMuPDF
                    logger.info("使用 PyMuPDF库提取PDF文本")
                    doc = fitz.open(file_path)
                    text = ""
                    for page_num in range(doc.page_count):
                        page = doc[page_num]
                        text += page.get_text()
                    doc.close()
                    
                    if text.strip():
                        logger.info(f"使用PyMuPDF成功提取PDF文本，长度: {len(text)}")
                    else:
                        logger.warning("使用PyMuPDF也未提取到任何文本")
                except ImportError:
                    logger.warning("PyMuPDF库也未安装，无法提取PDF文本")
                    return ""
                except Exception as e:
                    logger.error(f"使用PyMuPDF提取PDF文本失败: {str(e)}")
                    return ""
            except Exception as e:
                logger.error(f"使用PyPDF2提取PDF文本失败: {str(e)}")
                return ""
        except Exception as e:
            logger.error(f"使用pypdf提取PDF文本失败: {str(e)}")
            return ""
        
        if not text.strip():
            logger.warning(f"PDF文件 {file_path} 未能提取到任何文本内容")
        else:
            logger.info(f"成功从PDF文件提取 {len(text)} 个字符的文本")
    except Exception as e:
        logger.error(f"提取PDF文本失败: {str(e)}", exc_info=True)
        raise
    return text

def extract_text_from_docx(file_path: str) -> str:
    """从DOCX文件中提取文本"""
    text = ""
    try:
        logger.info(f"开始提取DOCX文件内容: {file_path}")
        from docx import Document
        doc = Document(file_path)
        paragraph_count = len(doc.paragraphs)
        logger.info(f"DOCX文件共有 {paragraph_count} 段落")
        
        for para_num, paragraph in enumerate(doc.paragraphs):
            try:
                if paragraph.text:
                    text += paragraph.text + "\n"
            except Exception as e:
                logger.error(f"提取第 {para_num+1} 段落文本失败: {str(e)}")
                continue
                
        if not text.strip():
            logger.warning(f"DOCX文件 {file_path} 未能提取到任何文本内容")
        else:
            logger.info(f"成功从DOCX文件提取 {len(text)} 个字符的文本")
    except Exception as e:
        logger.error(f"提取DOCX文本失败: {str(e)}", exc_info=True)
        raise
    return text

def extract_text_from_pptx(file_path: str) -> str:
    """从PPTX文件中提取文本"""
    text = ""
    try:
        logger.info(f"开始提取PPTX文件内容: {file_path}")
        from pptx import Presentation
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        logger.info(f"PPTX文件共有 {slide_count} 张幻灯片")
        
        for slide_num, slide in enumerate(prs.slides):
            try:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                if slide_text:
                    text += f"幻灯片 {slide_num+1}:\n{slide_text}\n"
                else:
                    logger.warning(f"第 {slide_num+1} 张幻灯片未能提取到文本")
            except Exception as e:
                logger.error(f"提取第 {slide_num+1} 张幻灯片文本失败: {str(e)}")
                continue
                
        if not text.strip():
            logger.warning(f"PPTX文件 {file_path} 未能提取到任何文本内容")
        else:
            logger.info(f"成功从PPTX文件提取 {len(text)} 个字符的文本")
    except Exception as e:
        logger.error(f"提取PPTX文本失败: {str(e)}", exc_info=True)
        raise
    return text

def extract_text_from_txt(file_path: str) -> str:
    """从TXT文件中提取文本"""
    try:
        logger.info(f"开始提取TXT文件内容: {file_path}")
        
        # 尝试使用UTF-8编码
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            logger.info(f"成功使用UTF-8编码读取TXT文件")
        except UnicodeDecodeError:
            # 尝试其他编码
            encodings = ['gbk', 'gb2312', 'gb18030', 'big5', 'latin1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        text = file.read()
                    logger.info(f"成功使用 {encoding} 编码读取TXT文件")
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise ValueError("无法确定文件编码")
                
        if not text.strip():
            logger.warning(f"TXT文件 {file_path} 未能提取到任何文本内容")
        else:
            logger.info(f"成功从TXT文件提取 {len(text)} 个字符的文本")
    except Exception as e:
        logger.error(f"提取TXT文本失败: {str(e)}", exc_info=True)
        raise
    return text

def extract_text_from_html(file_path: str) -> str:
    """从HTML文件中提取文本"""
    try:
        logger.info(f"开始提取HTML文件内容: {file_path}")
        from bs4 import BeautifulSoup
        with open(file_path, 'r', encoding='utf-8') as file:
            html_content = file.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        text = soup.get_text(separator='\n', strip=True)
        
        if not text.strip():
            logger.warning(f"HTML文件 {file_path} 未能提取到任何文本内容")
        else:
            logger.info(f"成功从HTML文件提取 {len(text)} 个字符的文本")
    except Exception as e:
        logger.error(f"提取HTML文本失败: {str(e)}", exc_info=True)
        raise
    return text

def extract_text_from_ipynb(file_path: str) -> str:
    """从IPYNB文件中提取文本"""
    text = ""
    try:
        logger.info(f"开始提取IPYNB文件内容: {file_path}")
        import nbformat
        with open(file_path, 'r', encoding='utf-8') as file:
            notebook = nbformat.read(file, as_version=4)
        
        cell_count = len(notebook.cells)
        logger.info(f"IPYNB文件共有 {cell_count} 个单元格")
        
        for cell_num, cell in enumerate(notebook.cells):
            try:
                if cell.cell_type in ['code', 'markdown']:
                    text += f"单元格 {cell_num+1} ({cell.cell_type}):\n{cell.source}\n\n"
            except Exception as e:
                logger.error(f"提取第 {cell_num+1} 个单元格文本失败: {str(e)}")
                continue
                
        if not text.strip():
            logger.warning(f"IPYNB文件 {file_path} 未能提取到任何文本内容")
        else:
            logger.info(f"成功从IPYNB文件提取 {len(text)} 个字符的文本")
    except Exception as e:
        logger.error(f"提取IPYNB文本失败: {str(e)}", exc_info=True)
        raise
    return text

# ==================== 文本分块模块 ====================

def split_text_by_chapters(text: str) -> List[Dict[str, Union[int, str]]]:
    """按章节分割文本（基于标题模式）"""
    try:
        # 定义常见的章节标题模式
        chapter_patterns = [
            r'^第[一二三四五六七八九十百千零\d]+章\s+.+$',  # 第一章 标题
            r'^第[一二三四五六七八九十百千零\d]+节\s+.+$',  # 第一节 标题
            r'^\d+\.\s+.+$',  # 1. 标题
            r'^\d+\.\d+\s+.+$',  # 1.1 标题
            r'^\d+\.\d+\.\d+\s+.+$',  # 1.1.1 标题
            r'^[A-Z][a-z]+\s+.+$',  # Title
            r'^[A-Z][a-z]+\s+[A-Z][a-z]+\s+.+$',  # Chapter Title
        ]
        
        # 合并所有模式
        combined_pattern = '|'.join(f'({pattern})' for pattern in chapter_patterns)
        
        # 查找所有章节标题
        chapter_matches = list(re.finditer(combined_pattern, text, re.MULTILINE))
        
        # 如果没有找到章节标题，返回整个文本作为一个块
        if not chapter_matches:
            return [{"title": "全文", "content": text}]
        
        # 分割文本
        chunks = []
        for i, match in enumerate(chapter_matches):
            title = match.group().strip()
            start = match.start()
            end = chapter_matches[i+1].start() if i+1 < len(chapter_matches) else len(text)
            content = text[start:end].strip()
            
            if content:  # 确保内容不为空
                chunks.append({
                    "title": title,
                    "content": content
                })
        
        return chunks
    except Exception as e:
        logger.error(f"按章节分割文本失败: {e}", exc_info=True)
        return [{"title": "全文", "content": text}]

def split_text_by_semantic(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """按语义分割文本（基于段落和句子）"""
    try:
        # 先按段落分割
        paragraphs = text.split('\n\n')
        
        # 初始化结果列表和当前块
        chunks = []
        current_chunk = ""
        
        for paragraph in paragraphs:
            # 如果当前段落比目标大小还大，按句子分割
            if len(paragraph) > chunk_size:
                # 先保存当前块（如果有内容）
                if current_chunk:
                    chunks.append(current_chunk)
                    current_chunk = ""
                
                # 按句子分割大段落
                sentences = re.split(r'(?<=[.!?。！？])\s+', paragraph)
                
                for sentence in sentences:
                    # 如果添加这个句子不会超过目标大小，就添加到当前块
                    if len(current_chunk) + len(sentence) <= chunk_size:
                        current_chunk += sentence + " "
                    else:
                        # 否则，保存当前块并开始新块
                        if current_chunk:
                            chunks.append(current_chunk)
                            current_chunk = sentence + " "
            else:
                # 如果添加这个段落不会超过目标大小，就添加到当前块
                if len(current_chunk) + len(paragraph) <= chunk_size:
                    current_chunk += paragraph + "\n\n"
                else:
                    # 否则，保存当前块并开始新块
                    if current_chunk:
                        chunks.append(current_chunk)
                        current_chunk = paragraph + "\n\n"
        
        # 添加最后一个块（如果有内容）
        if current_chunk:
            chunks.append(current_chunk)
        
        # 处理重叠
        if overlap > 0 and len(chunks) > 1:
            overlapping_chunks = []
            for i, chunk in enumerate(chunks):
                if i == 0:
                    overlapping_chunks.append(chunk)
                else:
                    # 获取前一个块的末尾部分
                    prev_chunk_end = chunks[i-1][-overlap:] if len(chunks[i-1]) > overlap else chunks[i-1]
                    # 添加到当前块的开头
                    overlapping_chunks.append(prev_chunk_end + chunk)
            chunks = overlapping_chunks
        
        return chunks
    except Exception as e:
        logger.error(f"按语义分割文本失败: {e}", exc_info=True)
        return [text]

def smart_split_text(text: str, file_type: str = "txt") -> List[Dict[str, Any]]:
    """智能分割文本（结合章节和语义分割）"""
    try:
        # 对于PDF等结构化文档，先按章节分割
        if file_type == "pdf":
            chapter_chunks = split_text_by_chapters(text)
            
            # 对每个章节内容进行语义分割
            result = []
            chunk_id = 0
            for chapter in chapter_chunks:
                title = chapter["title"]
                content = chapter["content"]
                
                # 如果章节内容较短，直接作为一个块
                if len(content) <= 1000:
                    result.append({
                        "chunk_id": chunk_id,
                        "title": title,
                        "content": content,
                        "metadata": {
                            "type": "chapter",
                            "title": title
                        }
                    })
                    chunk_id += 1
                else:
                    # 否则，对章节内容进行语义分割
                    semantic_chunks = split_text_by_semantic(content)
                    for i, semantic_chunk in enumerate(semantic_chunks):
                        result.append({
                            "chunk_id": chunk_id,
                            "title": f"{title} (部分 {i+1}/{len(semantic_chunks)})",
                            "content": semantic_chunk,
                            "metadata": {
                                "type": "chapter_part",
                                "title": title,
                                "part": i+1,
                                "total_parts": len(semantic_chunks)
                            }
                        })
                        chunk_id += 1
            
            return result
        else:
            # 对于其他类型文档，直接进行语义分割
            semantic_chunks = split_text_by_semantic(text)
            result = []
            for i, chunk in enumerate(semantic_chunks):
                result.append({
                    "chunk_id": i,
                    "title": f"文本块 {i+1}",
                    "content": chunk,
                    "metadata": {
                        "type": "semantic",
                        "chunk_index": i
                    }
                })
            
            return result
    except Exception as e:
        logger.error(f"智能分割文本失败: {e}", exc_info=True)
        return [{
            "chunk_id": 0,
            "title": "全文",
            "content": text,
            "metadata": {
                "type": "full_text"
            }
        }]

# ==================== MinIO 存储模块 ====================

def upload_file_stream_to_minio(
    file_content: bytes,
    object_name: str,
    file_type: str,
    bucket_name: str = None,
    metadata: Dict[str, Any] = None
) -> Tuple[bool, str]:
    """
    直接上传文件流到 MinIO
    
    Args:
        file_content: 文件内容（字节）
        object_name: 对象名称（在 MinIO 中的文件名）
        file_type: 文件类型（用于确定存储路径）
        bucket_name: 存储桶名称（可选，默认使用 Config.MINIO_BUCKET）
        metadata: 元数据（可选）
        
    Returns:
        (success, minio_path): 成功标志和 MinIO 对象路径
    """
    try:
        if bucket_name is None:
            bucket_name = Config.MINIO_BUCKET
        
        # 检查存储桶是否存在，不存在则创建
        found = minio_client.bucket_exists(bucket_name)
        if not found:
            minio_client.make_bucket(bucket_name)
            logger.info(f"创建 MinIO 存储桶: {bucket_name}")
        
        # 根据文件类型确定存储路径
        if file_type in ['pdf', 'docx', 'pptx', 'txt', 'html', 'ipynb']:
            object_path = f"{Config.MINIO_COURSE_MATERIALS_PATH}/{object_name}"
        else:
            object_path = f"{Config.MINIO_GENERAL_FILES_PATH}/{object_name}"
        
        # 上传文件流
        minio_client.put_object(
            bucket_name,
            object_path,
            data=io.BytesIO(file_content),
            length=len(file_content),
            metadata=metadata or {}
        )
        
        minio_path = f"minio://{bucket_name}/{object_path}"
        logger.info(f"文件流已成功上传到 MinIO: {minio_path}")
        
        return True, minio_path
        
    except S3Error as e:
        logger.error(f"MinIO S3 错误: {e}", exc_info=True)
        return False, f"MinIO S3 错误: {str(e)}"
    except Exception as e:
        logger.error(f"上传文件流到 MinIO 失败: {e}", exc_info=True)
        return False, f"上传文件流到 MinIO 失败: {str(e)}"

def download_file_from_minio(object_name: str, file_path: str, bucket_name: str = None) -> Tuple[bool, str]:
    """
    从 MinIO 下载文件
    
    Args:
        object_name: 对象名称（在 MinIO 中的文件名）
        file_path: 本地保存路径
        bucket_name: 存储桶名称（可选，默认使用 Config.MINIO_BUCKET）
        
    Returns:
        (success, downloaded_path): 成功标志和下载后的文件路径
    """
    try:
        if bucket_name is None:
            bucket_name = Config.MINIO_BUCKET
        
        # 构建对象路径（假设所有文件都放在 course_materials 目录下）
        object_path = f"{Config.MINIO_COURSE_MATERIALS_PATH}/{object_name}"
        
        # 下载文件
        minio_client.fget_object(
            bucket_name,
            object_path,
            file_path
        )
        
        logger.info(f"文件已从 MinIO 下载到: {file_path}")
        
        return True, file_path
        
    except S3Error as e:
        logger.error(f"MinIO S3 错误: {e}", exc_info=True)
        return False, f"MinIO S3 错误: {str(e)}"
    except Exception as e:
        logger.error(f"从 MinIO 下载文件失败: {e}", exc_info=True)
        return False, f"从 MinIO 下载文件失败: {str(e)}"

# ==================== 文件验证模块 ====================

def validate_file(file_path: str) -> Tuple[bool, Optional[str]]:
    """
    综合验证文件（存在性、大小、格式）
    """
    # 1. 检查文件是否存在
    if not os.path.exists(file_path):
        return False, f"❌ 文件不存在：{file_path}"
    
    # 2. 检查文件大小
    file_size = os.path.getsize(file_path)
    if file_size > Config.MAX_FILE_SIZE:
        size_mb = file_size / (1024 * 1024)
        return False, f"❌ 文件大小超过{Config.MAX_FILE_SIZE/(1024*1024)}MB限制，当前大小: {size_mb:.2f}MB"
    
    # 3. 检查文件格式
    file_extension = os.path.splitext(file_path)[1].lower().lstrip('.')
    if file_extension not in Config.ALLOWED_EXTENSIONS:
        return False, f"❌ 不支持的文件格式: {file_extension}，支持的格式: {', '.join(Config.ALLOWED_EXTENSIONS)}"
    
    return True, None

def get_supported_extensions() -> List[str]:
    """获取支持的文件扩展名列表"""
    return list(Config.ALLOWED_EXTENSIONS)

def get_max_file_size() -> int:
    """获取最大文件大小限制"""
    return Config.MAX_FILE_SIZE

def get_max_file_size_mb() -> float:
    """获取最大文件大小限制（MB）"""
    return Config.MAX_FILE_SIZE / (1024 * 1024)

def format_file_size(size_bytes: int) -> str:
    """格式化文件大小显示"""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.2f} {unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.2f} TB"

# ==================== 文件相关性检查模块（优化版）====================

def check_file_relevance(text_content: str, file_name: str) -> Dict[str, Any]:
    """
    检查文件内容是否与课程相关（优化版）
    
    Args:
        text_content: 文件文本内容
        file_name: 文件名
        
    Returns:
        相关性检查结果
    """
    # 1. 定义课程相关关键词
    course_keywords = [
        "课程", "课件", "知识点", "大纲", "习题", "作业", "考试", 
        "成绩", "教学", "学习", "章节", "内容", "资料", "教材", 
        "讲义", "授课", "教师", "教授", "讲师", "课堂", "教室", 
        "课时", "学分", "题目", "问题", "答案", "解析", "案例", 
        "实例", "示例", "练习", "实验", "实践", "项目", "设计", 
        "论文", "报告", "调研", "研究", "理论", "概念", "原理", 
        "方法", "技术", "应用", "发展", "趋势", "历史", "现状", 
        "未来", "前沿", "热点", "难点", "重点", "考点", "复习",
        "预习", "总结", "笔记", "教程", "指南", "入门", "基础",
        "进阶", "高级", "初级", "中级", "期末", "期中", "月考",
        "测验", "测试", "评估", "评价", "反馈", "建议", "指导",
        "教学计划", "教学大纲", "教学目标", "教学要求", "参考资料",
        "推荐阅读", "扩展阅读", "课外阅读", "辅助材料", "补充材料",
        "案例分析", "案例研究", "实战项目", "综合项目", "课程设计",
        "毕业设计", "课程论文", "课程报告", "课程作业", "课程实验",
        "课程实习", "课程实践", "课程考核", "课程评价", "课程评估",
        "课程反馈", "课程建议", "课程指导", "课程辅导", "课程答疑",
        "课程讨论", "课程交流", "课程分享", "课程资源", "课程工具",
        "课程平台", "课程系统", "课程网站", "课程APP", "课程小程序",
        "在线课程", "网络课程", "远程课程", "视频课程", "音频课程",
        "直播课程", "录播课程", "混合课程", "翻转课堂", "慕课",
        "微课", "公开课", "精品课", "示范课", "观摩课", "研讨课",
        "讲座", "报告会", "研讨会", "座谈会", "论坛", "沙龙"
    ]
    
    # 2. 定义非课程相关关键词（用于降低相关性评分）
    non_course_keywords = [
        "游戏", "娱乐", "小说", "电影", "音乐", "明星", "八卦",
        "购物", "美食", "旅游", "时尚", "美妆", "护肤", "减肥",
        "健身", "运动", "体育", "足球", "篮球", "网球", "游泳",
        "跑步", "瑜伽", "舞蹈", "唱歌", "乐器", "绘画", "摄影",
        "宠物", "动物", "植物", "花卉", "园艺", "烹饪", "美食",
        "汽车", "房产", "股票", "基金", "理财", "保险", "贷款",
        "信用卡", "借记卡", "支付宝", "微信支付", "银行", "证券",
        "期货", "外汇", "黄金", "白银", "原油", "天然气", "煤炭",
        "电力", "水利", "交通", "运输", "物流", "仓储", "供应链",
        "制造业", "建筑业", "房地产", "金融业", "互联网", "科技",
        "人工智能", "大数据", "云计算", "区块链", "物联网", "5G",
        "6G", "芯片", "半导体", "电子", "通信", "软件", "硬件",
        "操作系统", "数据库", "编程", "开发", "测试", "运维", "安全",
        "网络", "服务器", "云计算", "边缘计算", "雾计算", "区块链"
    ]
    
    # 3. 初始化相关性评分
    relevance_score = 0.0
    matched_keywords = []
    
    # 4. 检查课程相关关键词
    course_keyword_count = 0
    for keyword in course_keywords:
        if keyword in text_content:
            course_keyword_count += text_content.count(keyword)
            matched_keywords.append(keyword)
    
    # 5. 检查非课程相关关键词
    non_course_keyword_count = 0
    for keyword in non_course_keywords:
        if keyword in text_content:
            non_course_keyword_count += text_content.count(keyword)
    
    # 6. 计算相关性评分
    if course_keyword_count > 0:
        # 如果有课程相关关键词，基础分为 0.5
        base_score = 0.5
        # 根据关键词数量加分，最多加 0.5
        keyword_score = min(course_keyword_count / 10, 0.5)
        relevance_score = base_score + keyword_score
    else:
        # 如果没有课程相关关键词，基础分为 0.0
        base_score = 0.0
        relevance_score = base_score
    
    # 7. 如果有非课程相关关键词，降低相关性评分
    if non_course_keyword_count > 0:
        # 根据非课程相关关键词数量扣分，最多扣 0.3
        penalty = min(non_course_keyword_count / 10, 0.3)
        relevance_score = max(relevance_score - penalty, 0.0)
    
    # 8. 根据文件名调整相关性评分
    file_name_lower = file_name.lower()
    if any(keyword in file_name_lower for keyword in ["课程", "课件", "讲义", "大纲", "习题", "作业", "考试", "成绩"]):
        relevance_score = min(relevance_score + 0.1, 1.0)
    elif any(keyword in file_name_lower for keyword in ["游戏", "娱乐", "小说", "电影", "音乐"]):
        relevance_score = max(relevance_score - 0.2, 0.0)
    
    # 9. 根据文件扩展名调整相关性评分
    file_extension = os.path.splitext(file_name)[1].lower()
    if file_extension in ['pdf', 'docx', 'pptx', 'txt', 'html', 'ipynb']:
        # 如果是文档类文件，略微提高相关性评分
        relevance_score = min(relevance_score + 0.05, 1.0)
    elif file_extension in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'mp4', 'avi', 'mp3', 'wav']:
        # 如果是媒体类文件，略微降低相关性评分
        relevance_score = max(relevance_score - 0.1, 0.0)
    
    # 10. 判断是否相关
    is_relevant = relevance_score > 0.5
    
    # 11. 生成判断理由
    if is_relevant:
        if course_keyword_count > 0:
            reason = f"文件内容包含 {course_keyword_count} 个课程相关关键词：{', '.join(matched_keywords[:5])}"
        else:
            reason = "文件名或文件类型表明其可能是课程相关材料"
    else:
        if non_course_keyword_count > 0:
            reason = f"文件内容包含较多非课程相关关键词（{non_course_keyword_count} 个），相关性较低"
        else:
            reason = "文件内容未检测到足够的课程相关关键词"
    
    # 12. 生成建议课程类别
    suggested_categories = []
    if is_relevant:
        # 根据匹配的关键词生成建议课程类别
        if "编程" in matched_keywords or "代码" in matched_keywords or "开发" in matched_keywords:
            suggested_categories.append("计算机科学")
        if "数学" in matched_keywords or "微积分" in matched_keywords or "线性代数" in matched_keywords:
            suggested_categories.append("数学")
        if "物理" in matched_keywords or "力学" in matched_keywords or "电磁学" in matched_keywords:
            suggested_categories.append("物理")
        if "化学" in matched_keywords or "有机化学" in matched_keywords or "无机化学" in matched_keywords:
            suggested_categories.append("化学")
        if "英语" in matched_keywords or "语法" in matched_keywords or "词汇" in matched_keywords:
            suggested_categories.append("英语")
        if "历史" in matched_keywords or "近代史" in matched_keywords or "古代史" in matched_keywords:
            suggested_categories.append("历史")
        if "文学" in matched_keywords or "诗歌" in matched_keywords or "小说" in matched_keywords:
            suggested_categories.append("文学")
        
        # 如果没有匹配到具体类别，使用默认类别
        if not suggested_categories:
            suggested_categories = ["通识教育"]
    else:
        suggested_categories = ["其他"]
    
    # 13. 返回相关性检查结果
    return {
        "is_relevant": is_relevant,
        "relevance_score": relevance_score,
        "reason": reason,
        "suggested_categories": suggested_categories
    }

def handle_file_upload_success(file_info: Dict, current_thread_id: str = None) -> bool:
    """
    处理文件上传成功后的逻辑（优化版 - 包含图标系统和相关性检查）
    
    Args:
        file_info: 文件信息字典，包含文件名、大小、类型等
        current_thread_id: 当前会话ID（用于保存映射）
        
    Returns:
        是否处理成功
    """
    try:
        object_name = file_info.get("object_name")  # 带时间戳的文件名
        file_name = file_info.get("file_name")      # 原始文件名
        file_size = file_info.get("file_size")
        file_type = file_info.get("file_type")
        minio_path = file_info.get("minio_path")
        
        if current_thread_id is None:
            current_thread_id = "default"
        
        # 关键修复：保存原始文件名到带时间戳文件名的映射
        mapping_success = save_file_mapping(
            original_name=file_name,  # 原始文件名"成绩.txt"
            object_name=object_name,   # 带时间戳的文件名"成绩_1776341538.txt"
            thread_id=current_thread_id
        )
        
        if not mapping_success:
            logger.error(f"保存文件映射失败: {file_name} -> {object_name}")
        
        # 2. 检查文件相关性（仅当文件类型为文档类时）
        is_relevant = True
        relevance_score = 0.8
        reason = "文件内容包含大量课程相关关键词"
        suggested_categories = ["计算机科学", "编程"]
        
        if file_type in ['pdf', 'docx', 'pptx', 'txt', 'html', 'ipynb']:
            try:
                # 从 MinIO 下载文件到临时目录
                from .files_parser import download_file_from_minio
                from pathlib import Path
                
                temp_dir = Path(__file__).parent.parent / "data" / "temp"
                temp_dir.mkdir(parents=True, exist_ok=True)
                
                local_temp_path = temp_dir / object_name
                
                download_success, downloaded_path = download_file_from_minio(
                    object_name,
                    str(local_temp_path)
                )
                
                if download_success:
                    # 提取文件内容
                    text_content = extract_text_from_file(downloaded_path, file_type)
                    
                    if text_content:
                        # 调用相关性检查函数
                        relevance_result = check_file_relevance(text_content, file_name)
                        
                        is_relevant = relevance_result.get("is_relevant", True)
                        relevance_score = relevance_result.get("relevance_score", 0.8)
                        reason = relevance_result.get("reason", "文件内容包含大量课程相关关键词")
                        suggested_categories = relevance_result.get("suggested_categories", ["计算机科学", "编程"])
                        
                        logger.info(f"文件相关性检查结果: 相关={is_relevant}, 评分={relevance_score}, 理由={reason}")
                    else:
                        logger.warning(f"文件 {object_name} 内容为空或无法解析，默认认为相关")
                        is_relevant = True
                        relevance_score = 0.5
                        reason = "文件内容为空或无法解析，默认认为相关"
                        suggested_categories = ["通识教育"]
                else:
                    logger.warning(f"文件 {object_name} 下载失败，默认认为相关")
                    is_relevant = True
                    relevance_score = 0.5
                    reason = "文件下载失败，默认认为相关"
                    suggested_categories = ["通识教育"]
            except Exception as e:
                logger.error(f"检查文件相关性失败: {e}", exc_info=True)
                is_relevant = True
                relevance_score = 0.5
                reason = f"检查文件相关性失败: {str(e)}"
                suggested_categories = ["通识教育"]
        else:
            logger.info(f"文件类型 {file_type} 不是文档类，默认认为相关")
            is_relevant = True
            relevance_score = 0.8
            reason = f"文件类型 {file_type} 通常与课程相关"
            suggested_categories = ["通识教育"]
        
        # 3. 根据相关性决定是否更新向量库
        if is_relevant:
            # 将文件添加到向量库
            logger.info(f"开始将文件 {object_name} 添加到向量库")
            
            vector_success = add_documents_to_store_from_minio(object_name)
            
            if not vector_success:
                logger.error(f"将文件 {object_name} 添加到向量库失败")
                return False
        else:
            logger.info(f"文件 {object_name} 与课程相关性较低，不添加到向量库")
        
        # 4. 保存文件消息到聊天历史（优化版 - 使用新的图标系统）
        from .memory_sqlite import save_message
        
        # 生成文件卡片HTML（使用新的图标系统）
        file_card_html = generate_file_card_html(file_name, file_size, file_type, minio_path)
        
        # 生成状态卡片HTML
        if is_relevant:
            status_card_html = generate_file_status_card(
                file_name=file_name,
                status='success',
                message=f"文件已成功上传到向量库",
                file_extension=file_type
            )
        else:
            status_card_html = generate_file_status_card(
                file_name=file_name,
                status='warning',
                message=f"文件与课程相关性较低（评分: {relevance_score:.2f}），未上传到向量库",
                file_extension=file_type
            )
        
        # 保存文件卡片 (用户消息)
        save_message(
            thread_id=current_thread_id,
            role="user",
            content=file_card_html,
            message_type="file_card",
            metadata={
                "file_name": file_name,
                "file_size": file_size,
                "file_type": file_type,
                "object_name": object_name,
                "minio_path": minio_path,
                "upload_time": datetime.now().isoformat(),
                "is_relevant": is_relevant,
                "relevance_score": relevance_score,
                "reason": reason,
                "suggested_categories": suggested_categories
            }
        )
        
        # 保存状态卡片 (助手消息)
        save_message(
            thread_id=current_thread_id,
            role="assistant",
            content=status_card_html,
            message_type="file_status",
            metadata={
                "file_name": file_name,
                "file_size": file_size,
                "file_type": file_type,
                "object_name": object_name,
                "minio_path": minio_path,
                "upload_time": datetime.now().isoformat(),
                "is_relevant": is_relevant,
                "relevance_score": relevance_score,
                "reason": reason,
                "suggested_categories": suggested_categories,
                "status": "success" if is_relevant else "warning"
            }
        )
        
        logger.info(f"文件 {object_name} 处理完成")
        return True
        
    except Exception as e:
        logger.error(f"处理文件上传失败: {e}", exc_info=True)
        return False

# ==================== 测试模块 ====================

def test_file_parsing(file_path: str, file_extension: str = None) -> Dict[str, Any]:
    """
    测试文件解析功能
    
    Args:
        file_path: 要测试的文件路径
        file_extension: 文件扩展名（可选）
        
    Returns:
        包含测试结果的字典
    """
    result = {
        "file_path": file_path,
        "success": False,
        "error": None,
        "file_size": 0,
        "content_length": 0,
        "content_preview": "",
        "parsing_time": 0
    }
    
    import time
    start_time = time.time()
    
    try:
        # 检查文件是否存在
        if not os.path.exists(file_path):
            result["error"] = f"文件不存在: {file_path}"
            return result
        
        # 获取文件大小
        result["file_size"] = os.path.getsize(file_path)
        
        # 提取文件内容
        content = extract_text_from_file(file_path, file_extension)
        
        # 记录结果
        result["content_length"] = len(content)
        result["content_preview"] = content[:200] if content else ""
        result["success"] = bool(content)
        
        # 如果内容为空，记录警告
        if not content:
            result["error"] = "文件内容为空或无法解析"
        
    except Exception as e:
        result["error"] = str(e)
        result["success"] = False
    
    finally:
        result["parsing_time"] = time.time() - start_time
    
    return result

def test_file_from_minio(object_name: str, bucket_name: str = None) -> Dict[str, Any]:
    """
    测试从MinIO下载并解析文件
    
    Args:
        object_name: MinIO中的对象名
        bucket_name: 存储桶名称（可选）
        
    Returns:
        包含测试结果的字典
    """
    result = {
        "object_name": object_name,
        "success": False,
        "error": None,
        "download_success": False,
        "parsing_success": False,
        "download_time": 0,
        "parsing_time": 0,
        "content_length": 0,
        "content_preview": ""
    }
    
    import time
    from pathlib import Path
    
    try:
        # 创建临时目录
        temp_dir = Path(__file__).parent.parent / "data" / "temp"
        temp_dir.mkdir(parents=True, exist_ok=True)
        
        # 本地临时文件路径
        local_path = temp_dir / object_name
        
        # 1. 测试下载
        download_start = time.time()
        download_success, downloaded_path = download_file_from_minio(
            object_name,
            str(local_path),
            bucket_name
        )
        result["download_time"] = time.time() - download_start
        result["download_success"] = download_success
        
        if not download_success:
            result["error"] = f"下载失败: {downloaded_path}"
            return result
        
        # 2. 测试解析
        parsing_start = time.time()
        file_extension = os.path.splitext(object_name)[1].lower().lstrip('.')
        content = extract_text_from_file(downloaded_path, file_extension)
        result["parsing_time"] = time.time() - parsing_start
        result["parsing_success"] = bool(content)
        
        # 记录结果
        result["content_length"] = len(content)
        result["content_preview"] = content[:200] if content else ""
        result["success"] = download_success and result["parsing_success"]
        
        # 如果内容为空，记录警告
        if not content:
            result["error"] = "文件下载成功但内容为空或无法解析"
        
    except Exception as e:
        result["error"] = str(e)
        result["success"] = False
    
    return result

def display_test_result(result: Dict[str, Any]):
    """
    格式化显示测试结果
    
    Args:
        result: 测试结果字典
    """
    print("\n" + "="*60)
    print("文件解析测试结果")
    print("="*60)
    
    if result.get("file_path"):
        print(f"文件路径: {result['file_path']}")
    if result.get("object_name"):
        print(f"对象名称: {result['object_name']}")
    
    print(f"测试状态: {'✅ 成功' if result['success'] else '❌ 失败'}")
    
    if result.get("file_size"):
        print(f"文件大小: {format_file_size(result['file_size'])}")
    
    if result.get("download_success") is not None:
        print(f"下载状态: {'✅ 成功' if result['download_success'] else '❌ 失败'}")
        if result.get("download_time"):
            print(f"下载耗时: {result['download_time']:.2f}秒")
    
    if result.get("parsing_success") is not None:
        print(f"解析状态: {'✅ 成功' if result['parsing_success'] else '❌ 失败'}")
        if result.get("parsing_time"):
            print(f"解析耗时: {result['parsing_time']:.2f}秒")
    
    if result.get("content_length") is not None:
        print(f"内容长度: {result['content_length']} 字符")
    
    if result.get("content_preview"):
        print(f"\n内容预览:")
        print("-" * 60)
        print(result['content_preview'])
        print("-" * 60)
    
    if result.get("error"):
        print(f"\n错误信息: {result['error']}")
    
    print("="*60 + "\n")

# 命令行测试入口
if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="文件解析测试工具")
    parser.add_argument("--file", type=str, help="本地文件路径")
    parser.add_argument("--object", type=str, help="MinIO对象名称")
    parser.add_argument("--bucket", type=str, help="MinIO存储桶名称")
    
    args = parser.parse_args()
    
    if args.file:
        # 测试本地文件
        result = test_file_parsing(args.file)
        display_test_result(result)
    elif args.object:
        # 测试MinIO文件
        result = test_file_from_minio(args.object, args.bucket)
        display_test_result(result)
    else:
        print("请提供 --file 或 --object 参数")
        print("示例:")
        print("  测试本地文件: python files_parser.py --file /temp/成绩_1776341538.txt")
        print("  测试MinIO文件: python files_parser.py --object 成绩_1776341538.txt")
